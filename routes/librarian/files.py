from flask import Blueprint, render_template, request, redirect, url_for, jsonify
from routes.shared.auth import librarian_required
from db.organization_model import OrganizationModel
from db.conversation_model import ConversationModel
from db.file_model import FileModel
from db.user_model import UserModel
from docx import Document
from pptx import Presentation
import io

librarian = Blueprint('librarian', __name__)

@librarian.route('/librarian')
@librarian_required
def librarian_home(librarian):
    organization = OrganizationModel.get_organization(librarian['organization_id'])
    recent_conversations = ConversationModel.get_conversations_for_organization(librarian['organization_id'])
    total_prompt_history_count = ConversationModel.getTotalPromptHistoryEntriesForOrganization(librarian['organization_id'])
    conversation_and_question_count_for_all_users = ConversationModel.get_conversation_and_question_count_for_all_users(librarian['organization_id'])

    return render_template('librarian/librarian_home.html', librarian=librarian, organization=organization, recent_conversations=recent_conversations, total_prompt_history_count=total_prompt_history_count, conversation_and_question_count_for_all_users=conversation_and_question_count_for_all_users)

@librarian.route('/librarian/upload', methods=['POST'])
@librarian_required
def upload_file(librarian_user):
    if 'file' not in request.files:
        return redirect(url_for('librarian.librarian_files'))

    file = request.files['file']
    if file.filename == '':
        return redirect(url_for('librarian.librarian_files'))
    
    binary_content = file.read()
    filename = file.filename
    file_size = len(binary_content)
    
    # Detect file type and extract text content
    text_content = extract_text_content(file, binary_content)
    
    # Add the file to the files table
    FileModel.add_file(
        organization_id=librarian_user['organization_id'],
        user_upload_id=librarian_user['id'],
        binary_content=binary_content,
        text_content=text_content,
        file_name=filename,
        file_size=file_size
    )

    return redirect(url_for('librarian.librarian_files'))

def extract_text_content(file, binary_content):
    file_extension = file.filename.split('.')[-1].lower()
    
    if file_extension == 'docx':
        return extract_text_from_docx(binary_content)
    elif file_extension == 'pptx':
        return extract_text_from_pptx(binary_content)
    else:
        return decode_text_content(binary_content)

def extract_text_from_docx(binary_content):
    doc = Document(io.BytesIO(binary_content))
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    return "\n".join(text)

def extract_text_from_pptx(binary_content):
    prs = Presentation(io.BytesIO(binary_content))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def decode_text_content(binary_content):
    encodings = ['utf-8', 'iso-8859-1', 'windows-1252', 'ascii']
    for encoding in encodings:
        try:
            return binary_content.decode(encoding)
        except UnicodeDecodeError:
            print(f"Failed to decode with {encoding}")
            continue
    return None

@librarian.route('/librarian-files')
@librarian_required
def librarian_files(librarian_user):
    organization = OrganizationModel.get_organization(librarian_user['organization_id'])
    files = FileModel.get_files_for_organization(librarian_user['organization_id'])
    return render_template('librarian/librarian_files.html', librarian=librarian_user, organization=organization, files=files)

@librarian.route('/librarian/files', methods=['GET'])
@librarian_required
def get_organization_files(librarian_user):
    org_id = librarian_user['organization_id']
    
    files = FileModel.get_files_for_organization(org_id)
    users = UserModel.get_all_users()
    users_dict = {user['id']: user['username'] for user in users}
    print(f"users_dict: {users_dict}")
    
    file_list = []
    for file in files:
        file_list.append({
            'id': file['id'],
            'name': file['file_name'],
            'size': file['file_size'],
            'upload_date': file['timestamp_of_upload'],
            'uploaded_by': file['user_upload_id'],
            'user_name': users_dict[file['user_upload_id']]
        })
    
    return jsonify(file_list)

@librarian.route('/librarian/delete_file/<int:file_id>', methods=['POST'])
@librarian_required
def delete_file(librarian_user,file_id):
    file = FileModel.get_file_by_id(file_id)
    
    if file and file['organization_id'] == librarian_user['organization_id']:
        if FileModel.delete_file(file_id):
            return jsonify({'success': True, 'message': 'File deleted successfully'})
        else:
            return jsonify({'success': False, 'message': 'Failed to delete file'}), 500
    else:
        return jsonify({'success': False, 'message': 'File not found or you do not have permission to delete it'}), 404

@librarian.route('/librarian/preview_file/<int:file_id>', methods=['GET'])
@librarian_required
def preview_file(librarian_user, file_id):
    file = FileModel.get_file_content(file_id)
    if not file:
        return jsonify({"error": f"File with id {file_id} not found"}), 404

    return jsonify({"name": file['name'], "content": file['text_content']})
